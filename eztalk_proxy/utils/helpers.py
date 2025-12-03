import orjson
import re
import logging
import datetime
from typing import Any, Dict, List, Tuple, Optional
import os
import uuid

from fastapi.responses import JSONResponse
from fastapi import UploadFile, Depends
from fastapi.security import OAuth2PasswordBearer
from ..models.api_models import User
 
from ..core.config import (
    COMMON_HEADERS,
    MAX_SSE_LINE_LENGTH,
    SUPPORTED_DOCUMENT_MIME_TYPES_FOR_TEXT_EXTRACTION,
    MAX_DOCUMENT_CONTENT_CHARS_FOR_PROMPT,
)

try:
    from google.cloud import storage
    from google.auth.exceptions import DefaultCredentialsError
except ImportError:
    storage = None # type: ignore
    DefaultCredentialsError = None # type: ignore
    logging.warning(
        "google-cloud-storage or google-auth library not found. "
        "GCS upload functionality for large files (video/audio) for Gemini will not be available."
    )

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None
    logging.warning("1PyPDF2 library not found. PDF text extraction will not be available.")

try:
    import docx
except ImportError:
    docx = None
    logging.warning("python-docx library not found. DOCX text extraction will not be available.")

try:
    import olefile
except ImportError:
    olefile = None
    logging.warning("olefile library not found. DOC text extraction will not be available.")

try:
    import openpyxl
    import xlrd
except ImportError:
    openpyxl = None
    xlrd = None
    logging.warning("openpyxl/xlrd libraries not found. Excel text extraction will not be available.")

try:
    from pptx import Presentation
    pptx_available = True
except ImportError:
    pptx_available = False
    logging.warning("python-pptx library not found. PowerPoint text extraction will not be available.")

try:
    from bs4 import BeautifulSoup
    bs4_available = True
except ImportError:
    bs4_available = False
    logging.warning("BeautifulSoup4 library not found. HTML text extraction will be limited.")


logger = logging.getLogger("EzTalkProxy.Utils")

# 快速检测是否疑似包含 Markdown 结构，用于语音模式兜底清洗的早期退出
_MD_QUICK_PATTERN = re.compile(
    r"(^\s{0,3}#{1,6}\s)"          # 标题行: # / ## / ...
    r"|(^\s*[-*+]\s+)"            # 无序列表: - item
    r"|(^\s*\d+\.\s+)"            # 有序列表: 1. item
    r"|(```)"                     # 代码块围栏
    r"|(`[^`]+`)"                 # 行内代码
    r"|(\[.+?\]\(.+?\))"          # 链接 [text](url)
    r"|(!\[.*?\]\(.*?\))"         # 图片 ![alt](url)
    r"|(^\s*\|.+\|)"              # 表格行
    r"|(\*\*.+?\*\*)"             # 粗体 **text**
    ,
    re.MULTILINE | re.DOTALL,
)


def strip_markdown_for_tts(text: str) -> str:
    """
    语音模式兜底：在送入 TTS 前对明显的 Markdown 结构做轻量清洗。
    设计原则：
    - 检测与清洗都必须是 O(n) 且非常快
    - 若文本本身没有明显 Markdown 结构，则直接原样返回，不做任何修改
    - 尽量删除/改写排版符号，保留语义内容
    """
    if not isinstance(text, str) or not text:
        return text

    # 快速检测：绝大多数正常口语在这里直接返回，避免多轮正则
    if not _MD_QUICK_PATTERN.search(text):
        return text

    cleaned = text

    # 1. 去掉围栏代码块 ``` ```（整体移除，避免读代码噪音）
    cleaned = re.sub(
        r"```[a-zA-Z0-9_+\-]*\n.*?\n```",
        "",
        cleaned,
        flags=re.DOTALL,
    )

    # 2. 行首标题、引用、列表标记
    cleaned = re.sub(r"^\s{0,3}#{1,6}\s*", "", cleaned, flags=re.MULTILINE)
    cleaned = re.sub(r"^\s{0,3}>\s?", "", cleaned, flags=re.MULTILINE)
    cleaned = re.sub(r"^\s*[-*+]\s+", "", cleaned, flags=re.MULTILINE)
    cleaned = re.sub(r"^\s*\d+\.\s+", "", cleaned, flags=re.MULTILINE)

    # 3. 水平分隔线（---、*** 等）
    cleaned = re.sub(
        r"^\s{0,3}[-*_]{3,}\s*$",
        "",
        cleaned,
        flags=re.MULTILINE,
    )

    # 4. 链接与图片：[text](url)、![alt](url) -> 仅保留文字部分
    cleaned = re.sub(
        r"!\[([^\]]*)\]\([^)]+\)",
        r"\1",
        cleaned,
    )
    cleaned = re.sub(
        r"\[([^\]]+)\]\([^)]+\)",
        r"\1",
        cleaned,
    )

    # 5. 行内代码与加粗/斜体：去掉符号，保留内容
    cleaned = re.sub(r"`([^`]+)`", r"\1", cleaned)
    cleaned = re.sub(r"\*\*([^*]+)\*\*", r"\1", cleaned)
    cleaned = re.sub(r"\*([^*]+)\*", r"\1", cleaned)

    # 6. 表格中竖线：将 | 替换为顿号/逗号，避免被读出为奇怪符号
    # 这里不过度区分是不是表格，只要出现大量 | 就做温和替换
    cleaned = re.sub(r"\|+", "，", cleaned)

    # 7. 收尾：清理多余空白与空行
    cleaned = re.sub(r"[ \t]+\n", "\n", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    cleaned = cleaned.strip()

    return cleaned


def orjson_dumps_bytes_wrapper(data: Any) -> bytes:
    return orjson.dumps(
        data,
        option=orjson.OPT_NON_STR_KEYS | orjson.OPT_PASSTHROUGH_DATETIME | orjson.OPT_APPEND_NEWLINE
    )

def to_sse_bytes(event) -> bytes:
    """
    将 AppStreamEventPy 序列化为符合 SSE 规范的字节：
    - 前缀使用 "data: "（注意后面的空格，兼容性更好）
    - 每条事件以双换行结尾 \\n\\n
    - 允许输入已带 "data:" 的载荷（幂等处理）
    """
    try:
        # 统一生成 JSON bytes（末尾通常带一个 \\n）
        payload = orjson_dumps_bytes_wrapper(
            event.model_dump(by_alias=True, exclude_none=True)
        )
        # 如果上游已带 "data:" 前缀，保留；否则补充 "data: "
        if payload.startswith(b"data:"):
            line = payload
        else:
            line = b"data: " + payload

        # 确保以两个换行结束（SSE事件分隔）
        if line.endswith(b"\n\n"):
            return line
        elif line.endswith(b"\n"):
            return line + b"\n"
        else:
            return line + b"\n\n"
    except Exception:  # 兜底：任何异常都返回一个可用的 error 事件，避免客户端挂起
        fallback = {
            "type": "error",
            "message": "SSE serialization failure",
        }
        fb = orjson_dumps_bytes_wrapper(fallback)
        if not fb.startswith(b"data:"):
            fb = b"data: " + fb
        if fb.endswith(b"\n\n"):
            return fb
        elif fb.endswith(b"\n"):
            return fb + b"\n"
        else:
            return fb + b"\n\n"
def error_response(
    code: int,
    msg: str,
    request_id: Optional[str] = None,
    headers: Optional[Dict[str, str]] = None
) -> JSONResponse:
    log_msg = f"错误 {code}: {msg}"
    if request_id:
        log_msg = f"RID-{request_id}: {log_msg}"
    logger.warning(log_msg)
    
    final_headers = {**COMMON_HEADERS, **(headers or {})}
    
    return JSONResponse(
        status_code=code,
        content={"error": {"message": msg, "code": code, "type": "proxy_error"}},
        headers=final_headers
    )


def extract_sse_lines(buffer: bytearray) -> Tuple[List[bytes], bytearray]:
    lines: List[bytes] = []
    start_index: int = 0
    buffer_len = len(buffer)
    while start_index < buffer_len:
        newline_index = buffer.find(b'\n', start_index)
        if newline_index == -1:
            break
        line = buffer[start_index:newline_index]
        if line.endswith(b'\r'):
            line = line[:-1]
        if len(line) > MAX_SSE_LINE_LENGTH:
            logger.warning(
                f"SSE line too long ({len(line)} bytes), exceeded MAX_SSE_LINE_LENGTH ({MAX_SSE_LINE_LENGTH}). Line skipped. "
                f"Content start: {line[:100]!r}"
            )
        else:
            lines.append(line)
        start_index = newline_index + 1
    return lines, buffer[start_index:]

def get_current_time_iso() -> str:
    return datetime.datetime.utcnow().isoformat() + "Z"

def is_gemini_2_5_model(model_name: str) -> bool:
    if not isinstance(model_name, str):
        return False
    return "gemini-2.5" in model_name.lower() or "gemini-2.5-flash-image-preview" in model_name.lower()

def _extract_text_from_pdf_pypdf2(file_path: str) -> Optional[str]:
    if not PyPDF2:
        logger.warning("Attempted to extract PDF text, but PyPDF2 library is not available.")
        return None
    text_content = ""
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            if reader.is_encrypted:
                try:
                    if reader.decrypt("") == PyPDF2.PasswordType.OWNER_PASSWORD or \
                       reader.decrypt("") == PyPDF2.PasswordType.USER_PASSWORD :
                        logger.info(f"Successfully decrypted PDF (with empty password): {file_path}")
                    else:
                        logger.warning(f"PDF file is encrypted and could not be decrypted with an empty password: {file_path}")
                        return None
                except Exception as decrypt_err:
                    logger.warning(f"Failed to decrypt PDF {file_path}: {decrypt_err}")
                    return None

            for page in reader.pages:
                try:
                    text_content += page.extract_text() or ""
                except Exception as page_extract_err:
                    logger.warning(f"Error extracting text from a page in {file_path}: {page_extract_err}")
                    continue
        return text_content.strip()
    except FileNotFoundError:
        logger.error(f"PDF file not found for extraction: {file_path}")
        return None
    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_path} using PyPDF2: {e}", exc_info=True)
        return None

def _extract_text_from_docx_python_docx(file_path: str) -> Optional[str]:
    if not docx:
        logger.warning("Attempted to extract DOCX text, but python-docx library is not available.")
        return None
    try:
        doc_obj = docx.Document(file_path)
        full_text = [para.text for para in doc_obj.paragraphs]
        return "\n".join(full_text).strip()
    except FileNotFoundError:
        logger.error(f"DOCX file not found for extraction: {file_path}")
        return None
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path} using python-docx: {e}", exc_info=True)
        return None

def _extract_text_from_doc_olefile(file_path: str) -> Optional[str]:
    """使用olefile库从.doc文档中提取文本（增强版，支持中文，智能编码检测）"""
    if not olefile:
        logger.warning("Attempted to extract DOC text, but olefile library is not available.")
        return None
    
    content = None
    try:
        # 尝试通过 OLE 结构读取 WordDocument 流，减少二进制噪声
        if olefile.isOleFile(file_path):
            with olefile.OleFileIO(file_path) as ole:
                if ole.exists('WordDocument'):
                    with ole.openstream('WordDocument') as stream:
                        content = stream.read()
    except Exception as e:
        logger.warning(f"Failed to read WordDocument stream from {file_path}: {e}")
    
    # 如果读取流失败，回退到读取整个文件
    if content is None:
        try:
            with open(file_path, 'rb') as f:
                content = f.read()
        except Exception as e:
            logger.error(f"Failed to read file {file_path}: {e}")
            return None

    # 辅助函数：评估文本看起来是否像正常的中文/英文文本
    def score_text_validity(text: str) -> float:
        if not text: return 0.0
        length = len(text)
        if length == 0: return 0.0
        
        # 常用中文高频字和标点
        common_chars = {'的', '一', '是', '在', '不', '了', '有', '和', '人', '这', '中', '大', '为', '上', '个', '国', '我', '以', '要', '他', '时', '来', '用', '们', '生', '到', '作', '地', '于', '出', '就', '分', '对', '成', '会', '可', '主', '发', '年', '动', '同', '工', '也', '能', '下', '过', '子', '说', '产', '种', '面', '而', '方', '后', '多', '定', '行', '学', '法', '所', '民', '得', '经', '十', '三', '之', '进', '着', '等', '部', '度', '家', '电', '力', '里', '如', '水', '化', '高', '自', '二', '理', '起', '小', '物', '现', '实', '加', '量', '都', '两', '体', '制', '机', '当', '使', '点', '从业', '本', '去', '心', '界', '义', '社', '合', '平', '士', '告', '外', '没', '看', '提', '那', '问', '指', '气', '做', '邻', '西', '真', '山', '内', '月', '公', '全', '信', '期', '安', '或', '书', '门', '应', '路', '利', '手', '最', '新', '世', '位', '场', '变', '得', '员', '表', '口', '常', '关', '争', '军', '目', '者', '次', '解', '文', '九', '八', '无', '相', '日', '外', '刚', '但', '步', '名', '建', '果', '料', '张', '接', '员', '司', '住', '实', '运', '通', '农', '保', '导', '集', '物', '展', '象', '完', '院', '样', '干', '并', '利', '省', '源', '安', '千', '众', '效', '管', '接', '觉', '身', '美', '意', '先', '金', '月', '回', '工', '热', '性', '音', '老', '切', '级', '由', '因', '联', '即', '百', '知', '表', '队', '组', '决', '治', '看', '住', '美', '点', '题', '，', '。', '、', '；', '：', '？', '！', '“', '”', '（', '）'}
        
        score = 0
        for char in text:
            if char in common_chars:
                score += 1
            elif 'a' <= char <= 'z' or 'A' <= char <= 'Z' or '0' <= char <= '9':
                score += 0.1 # 英文数字权重低一点，避免全英文二进制噪声干扰
        
        return score / length

    candidates = []
    import re

    # 方法1：UTF-16LE (Word默认)
    try:
        # 填充到偶数长度
        content_le = content + b'\0' if len(content) % 2 != 0 else content
        # 解码并过滤无效字符
        text_le = content_le.decode('utf-16-le', errors='ignore')
        # 清洗：只保留 CJK、ASCII 和常见符号
        # 匹配连续的有效字符块（至少2个字符），减少单字噪声
        matches = re.findall(r'[\u4e00-\u9fff\x20-\x7e\uff00-\uffef\u3000-\u303f\t\n\r]{2,}', text_le)
        cleaned_le = "".join(matches)
        score_le = score_text_validity(cleaned_le)
        candidates.append((score_le, cleaned_le, "utf-16-le"))
    except Exception:
        pass

    # 方法2：GB18030 (中文文档常见，尤其是老文档)
    try:
        text_gb = content.decode('gb18030', errors='ignore')
        matches = re.findall(r'[\u4e00-\u9fff\x20-\x7e\uff00-\uffef\u3000-\u303f\t\n\r]{2,}', text_gb)
        cleaned_gb = "".join(matches)
        score_gb = score_text_validity(cleaned_gb)
        candidates.append((score_gb, cleaned_gb, "gb18030"))
    except Exception:
        pass

    # 选取得分最高的
    candidates.sort(key=lambda x: x[0], reverse=True)
    
    if not candidates:
        return None

    best_score, best_text, best_encoding = candidates[0]
    
    # 阈值：至少有1%的字符是常见字符，否则认为是乱码
    if best_score > 0.01 and len(best_text) > 5:
        logger.info(f"Extracted .doc text using {best_encoding}, score={best_score:.4f}, length={len(best_text)}")
        return best_text
    
    logger.warning(f"Extraction failed: best score {best_score:.4f} too low. Encoding: {best_encoding}")
    return None

def _extract_text_from_excel(file_path: str, mime_type: str) -> Optional[str]:
    """从Excel文件中提取文本"""
    try:
        text_content = []
        
        if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            # .xlsx 文件
            if not openpyxl:
                logger.warning("openpyxl not available for .xlsx extraction")
                return None
                
            workbook = openpyxl.load_workbook(file_path)
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"=== Sheet: {sheet_name} ===")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        text_content.append(" | ".join(row_text))
                        
        elif mime_type == "application/vnd.ms-excel":
            # .xls 文件
            if not xlrd:
                logger.warning("xlrd not available for .xls extraction")
                return None
                
            workbook = xlrd.open_workbook(file_path)
            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                text_content.append(f"=== Sheet: {sheet.name} ===")
                
                for row_idx in range(sheet.nrows):
                    row_text = []
                    for col_idx in range(sheet.ncols):
                        cell = sheet.cell(row_idx, col_idx)
                        if cell.value:
                            row_text.append(str(cell.value))
                    if row_text:
                        text_content.append(" | ".join(row_text))
        
        return "\n".join(text_content).strip() if text_content else None
        
    except Exception as e:
        logger.error(f"Error extracting text from Excel {file_path}: {e}", exc_info=True)
        return None

def _extract_text_from_powerpoint(file_path: str) -> Optional[str]:
    """从PowerPoint文件中提取文本"""
    if not pptx_available:
        logger.warning("python-pptx not available for PowerPoint extraction")
        return None
        
    try:
        presentation = Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            text_content.append(f"=== Slide {slide_num} ===")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text.strip())
                    
        return "\n".join(text_content).strip() if text_content else None
        
    except Exception as e:
        logger.error(f"Error extracting text from PowerPoint {file_path}: {e}", exc_info=True)
        return None

def _extract_text_from_html(file_path: str) -> Optional[str]:
    """从HTML文件中提取文本"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
            
        if bs4_available:
            # 使用BeautifulSoup解析HTML
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # 移除script和style标签
            for script in soup(["script", "style"]):
                script.decompose()
                
            # 提取文本
            text = soup.get_text()
            
            # 清理多余的空白
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            return text.strip() if text else None
        else:
            # 简单的HTML标签移除
            import re
            # 移除HTML标签
            text = re.sub(r'<[^>]+>', '', html_content)
            # 清理多余空白
            text = re.sub(r'\s+', ' ', text)
            return text.strip() if text else None
            
    except Exception as e:
        logger.error(f"Error extracting text from HTML {file_path}: {e}", exc_info=True)
        return None

def _extract_text_from_xml(file_path: str) -> Optional[str]:
    """从XML文件中提取文本"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            xml_content = f.read()
            
        if bs4_available:
            soup = BeautifulSoup(xml_content, 'xml')
            text = soup.get_text()
            # 清理多余的空白
            lines = (line.strip() for line in text.splitlines())
            text = '\n'.join(line for line in lines if line)
            return text.strip() if text else None
        else:
            # 简单的XML标签移除
            import re
            text = re.sub(r'<[^>]+>', '', xml_content)
            text = re.sub(r'\s+', ' ', text)
            return text.strip() if text else None
            
    except Exception as e:
        logger.error(f"Error extracting text from XML {file_path}: {e}", exc_info=True)
        return None

def _extract_text_from_json(file_path: str) -> Optional[str]:
    """从JSON文件中提取文本内容"""
    try:
        import json
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            data = json.load(f)
            
        def extract_strings(obj, path=""):
            """递归提取JSON中的字符串值"""
            strings = []
            if isinstance(obj, dict):
                for key, value in obj.items():
                    new_path = f"{path}.{key}" if path else key
                    strings.extend(extract_strings(value, new_path))
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    new_path = f"{path}[{i}]"
                    strings.extend(extract_strings(item, new_path))
            elif isinstance(obj, str) and obj.strip():
                strings.append(f"{path}: {obj}")
            elif obj is not None:
                strings.append(f"{path}: {str(obj)}")
            return strings
            
        text_parts = extract_strings(data)
        return "\n".join(text_parts) if text_parts else None
        
    except Exception as e:
        logger.error(f"Error extracting text from JSON {file_path}: {e}", exc_info=True)
        return None

def _extract_text_from_plain_text(file_path: str) -> Optional[str]:
    common_encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1', 'iso-8859-1']
    try:
        for encoding in common_encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read().strip()
            except UnicodeDecodeError:
                logger.debug(f"Failed to decode plain text file {file_path} with encoding {encoding}")
                continue
            except FileNotFoundError:
                logger.error(f"Plain text file not found for extraction: {file_path}")
                return None
        logger.warning(f"Could not decode plain text file {file_path} with common encodings.")
        return None
    except Exception as e:
        logger.error(f"Error extracting text from plain text file {file_path}: {e}", exc_info=True)
        return None


async def extract_text_from_uploaded_document(
    uploaded_file_path: str,
    mime_type: Optional[str],
    original_filename: str
) -> Optional[str]:
    logger.info(f"Attempting to extract text from '{original_filename}' (path: {uploaded_file_path}, mime: {mime_type})")
    effective_mime_type = mime_type.lower() if mime_type else None

    if not effective_mime_type:
        logger.warning(f"No effective MIME type for '{original_filename}', cannot determine extraction method.")
        return None

    extracted_text: Optional[str] = None

    if effective_mime_type in SUPPORTED_DOCUMENT_MIME_TYPES_FOR_TEXT_EXTRACTION:
        # Microsoft Office Documents
        if effective_mime_type == "application/pdf":
            extracted_text = _extract_text_from_pdf_pypdf2(uploaded_file_path)
        elif effective_mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            extracted_text = _extract_text_from_docx_python_docx(uploaded_file_path)
        elif effective_mime_type == "application/msword":
            logger.warning(f"🔥 .doc格式文档处理：'{original_filename}' - .doc格式较老，提取效果可能不佳")
            extracted_text = _extract_text_from_doc_olefile(uploaded_file_path)
            if not extracted_text or len(extracted_text.strip()) < 10:
                extracted_text = f"""[文档解析提示]

.doc格式文档 '{original_filename}' 的内容提取遇到困难。

可能原因：
1. .doc是较老的Microsoft Word格式，结构复杂
2. 文档可能包含特殊格式或加密保护
3. 当前解析器对复杂.doc文档支持有限

建议解决方案：
1. 将文档转换为.docx格式后重新上传
2. 将文档另存为PDF格式后重新上传  
3. 复制文档内容到纯文本文件(.txt)后上传

如需帮助转换文档格式，请告知具体需求。"""
                logger.warning(f"Failed to extract meaningful content from .doc file '{original_filename}'")
        
        # Excel Documents
        elif effective_mime_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"]:
            extracted_text = _extract_text_from_excel(uploaded_file_path, effective_mime_type)
        
        # PowerPoint Documents  
        elif effective_mime_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
            extracted_text = _extract_text_from_powerpoint(uploaded_file_path)
        
        # Web & Markup Documents
        elif effective_mime_type == "text/html":
            extracted_text = _extract_text_from_html(uploaded_file_path)
        elif effective_mime_type in ["text/xml", "application/xml"]:
            extracted_text = _extract_text_from_xml(uploaded_file_path)
        elif effective_mime_type == "application/json":
            extracted_text = _extract_text_from_json(uploaded_file_path)
        
        # Plain Text & Other Formats
        elif effective_mime_type.startswith("text/"):
            extracted_text = _extract_text_from_plain_text(uploaded_file_path)
        else:
            logger.info(f"MIME type '{effective_mime_type}' for '{original_filename}' is in supported list but no specific extractor implemented, attempting plain text.")
            extracted_text = _extract_text_from_plain_text(uploaded_file_path)
    else:
        logger.warning(f"Unsupported MIME type for text extraction: '{effective_mime_type}' for file '{original_filename}'.")
        return None

    if extracted_text:
        if len(extracted_text) > MAX_DOCUMENT_CONTENT_CHARS_FOR_PROMPT:
            logger.info(f"Extracted text from '{original_filename}' truncated from {len(extracted_text)} to {MAX_DOCUMENT_CONTENT_CHARS_FOR_PROMPT} characters.")
            extracted_text = extracted_text[:MAX_DOCUMENT_CONTENT_CHARS_FOR_PROMPT] + \
                             f"\n[内容已截断，原始长度超过 {MAX_DOCUMENT_CONTENT_CHARS_FOR_PROMPT} 字符]"
        logger.info(f"Successfully extracted text (len: {len(extracted_text)}) from '{original_filename}'.")
        return extracted_text.strip()
    else:
        logger.warning(f"Failed to extract text from '{original_filename}' (mime: {effective_mime_type}).")
        return None

async def upload_to_gcs(
    file_obj: Any,
    original_filename: str,
    bucket_name: str,
    project_id: Optional[str] = None,
    content_type: Optional[str] = None,
    request_id: Optional[str] = None
) -> Optional[str]:
    log_prefix = f"RID-{request_id}" if request_id else "[GCS_UPLOAD]"
    
    if not storage:
        logger.error(f"{log_prefix} GCS upload skipped: google-cloud-storage library not available.")
        return None
    if not bucket_name:
        logger.error(f"{log_prefix} GCS upload skipped: GCS_BUCKET_NAME is not configured.")
        return None

    _, file_extension = os.path.splitext(original_filename)
    safe_original_filename_part = "".join(c if c.isalnum() or c in ['.', '_', '-'] else '_' for c in original_filename.rsplit('.', 1)[0])[:50]
    destination_blob_name = f"uploads/{request_id or 'unknown_req'}/{safe_original_filename_part}_{uuid.uuid4().hex[:8]}{file_extension}"

    logger.info(f"{log_prefix} Attempting to upload to GCS: bucket='{bucket_name}', blob='{destination_blob_name}'")

    try:
        if project_id:
            storage_client = storage.Client(project=project_id)
        else:
            storage_client = storage.Client()
            
        bucket = storage_client.bucket(bucket_name)
        blob = bucket.blob(destination_blob_name)

        if isinstance(file_obj, UploadFile):
            await file_obj.seek(0)
            blob.upload_from_file(file_obj.file, content_type=content_type or file_obj.content_type)
        elif hasattr(file_obj, 'read') and hasattr(file_obj, 'seek'):
            file_obj.seek(0)
            blob.upload_from_file(file_obj, content_type=content_type)
        elif isinstance(file_obj, str) and os.path.exists(file_obj):
            blob.upload_from_filename(file_obj, content_type=content_type)
        else:
            logger.error(f"{log_prefix} GCS upload failed for '{original_filename}': Invalid file_obj type or file path does not exist.")
            return None

        gcs_uri = f"gs://{bucket_name}/{destination_blob_name}"
        logger.info(f"{log_prefix} Successfully uploaded '{original_filename}' to GCS: {gcs_uri}")
        return gcs_uri
    except DefaultCredentialsError:
        logger.error(
            f"{log_prefix} GCS upload failed for '{original_filename}': Google Cloud Default Credentials not found. "
            "Ensure GOOGLE_APPLICATION_CREDENTIALS environment variable is set correctly "
            "or the runtime environment has appropriate GCS permissions.",
            exc_info=True
        )
        return None
    except Exception as e:
        logger.error(f"{log_prefix} GCS upload failed for '{original_filename}' (blob '{destination_blob_name}'): {e}", exc_info=True)
        return None

oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

async def get_current_user(token: str = Depends(oauth2_scheme)):
   # In a real application, you would verify the token and fetch the user from a database.
   # For now, we'll just return a dummy user.
   return User(username="johndoe", email="johndoe@example.com", full_name="John Doe")